import requests
import pandas as pd
from bs4 import BeautifulSoup

# Initialize variables
page = 1
all_reviews = []  # Stores all review data across pages

# Loop to iterate over multiple pages
for _ in range(389):
    data_list = []   # To hold the ratings for each review
    review_list = [] # To hold the review texts

    airline_URL = f"https://www.airlinequality.com/airline-reviews/british-airways/page/{page}/"
    response = requests.get(airline_URL)
    soup = BeautifulSoup(response.text, "html.parser")

    # Scrape reviews
    reviews = soup.find_all(name="div", class_="text_content")
    for review in reviews:
        text = review.get_text().strip()
        review_list.append(text)

    # Scrape ratings tables
    tables = soup.find_all(name="table", class_="review-ratings")
    for table in tables:
        data = {}
        rows = table.find_all("tr")
        for row in rows:
            header = row.find('td', class_='review-rating-header').text.strip()
            value = row.find('td', class_='review-value')
            rating_stars = row.find('td', class_='review-rating-stars')
            if value:
                data[header] = value.text.strip()
            elif rating_stars:
                stars_filled = rating_stars.find_all('span', class_='star fill')
                data[header] = len(stars_filled)

        data_list.append(data)

    data_list = data_list[1:]

    # Combine reviews and ratings
    for i in range(min(len(review_list), len(data_list))):
        review_data = {'review_text': review_list[i]}
        review_data.update(data_list[i])  # Add the rating data
        all_reviews.append(review_data)   # Add combined data to all_reviews

    print(f"Scraped page {page}")
    page += 1

# Convert to DataFrame and save to CSV
df = pd.DataFrame(all_reviews)
df.to_csv('data/british_airways_reviews.csv', index=False)
print("Data saved to 'data/british_airways_reviews.csv'")


df = pd.read_csv('data/british_airways_reviews.csv')

from textblob import TextBlob

def get_sentiment(n_text):
    return TextBlob(n_text).sentiment.polarity  # Returns a polarity score from -1 (negative) to +1 (positive)

df['sentiment'] = df['review_text'].apply(get_sentiment)

print(df)
print("Average Sentiment Score:", df['sentiment'].mean())

from wordcloud import WordCloud
import matplotlib.pyplot as plt
import seaborn as sns

word_cloud = WordCloud(width=800, height=400).generate(' '.join(df['review_text']))
plt.figure(figsize=(10, 5))
plt.imshow(word_cloud, interpolation='bilinear')
plt.axis('off')
plt.title("Word Cloud of Customer Reviews")
plt.savefig('word_cloud.png')
plt.show()

df['sentiment'].plot(kind='hist', bins=20, title='Sentiment Distribution')
plt.xlabel('Sentiment Score')
plt.show()

rating_columns = ['Seat Comfort', 'Cabin Staff Service', 'Food & Beverages',
                  'Ground Service', 'Value For Money', 'Inflight Entertainment', 'Wifi & Connectivity']

average_ratings = df[rating_columns].mean()
print(average_ratings)

# Plot the average ratings as a bar chart
plt.figure(figsize=(10, 6))
average_ratings.plot(kind='bar', color='skyblue')
plt.title('Average Ratings for British Airways Services')
plt.xlabel('Service Category')
plt.ylabel('Average Rating')
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig('average_ratings.png')
plt.show()

# Plot histograms for each rating category
plt.figure(figsize=(15, 10))
for i, col in enumerate(rating_columns, 1):
    plt.subplot(3, 3, i)
    sns.histplot(df[col], bins=5, kde=True, color='skyblue')
    plt.title(f'Distribution of {col}')
    plt.xlabel('Rating')
    plt.ylabel('Frequency')

plt.tight_layout()
plt.savefig('sentiment_distribution.png')
plt.show()

# Convert "Recommended" column to 1 (Yes) and 0 (No)
df['Recommended'] = df['Recommended'].apply(lambda x: 1 if x == 'yes' else 0)

# Calculate recommendation rate
recommend_rate = df['Recommended'].mean() * 100
print(f"Percentage of Customers Who Would Recommend the Airline: {recommend_rate:.2f}%")

# Plot pie chart of recommendation
recommend_counts = df['Recommended'].value_counts()
plt.figure(figsize=(6, 6))
plt.pie(recommend_counts, labels=['No', 'Yes'], autopct='%1.1f%%', startangle=140, colors=['coral', 'skyblue'])
plt.title('Customer Recommendation Rate')
plt.savefig('recommendation_rate.png')
plt.show()

from textblob import TextBlob

# Define a function to calculate sentiment
def get_sentiment(n_text):
    return TextBlob(n_text).sentiment.polarity

# Apply sentiment analysis to review_text
df['Sentiment Score'] = df['review_text'].apply(get_sentiment)

# Display average sentiment
average_sentiment = df['Sentiment Score'].mean()
print(f"Average Sentiment Score: {average_sentiment:.2f}")

# Plot sentiment score distribution
plt.figure(figsize=(10, 6))
sns.histplot(df['Sentiment Score'], bins=20, kde=True, color='skyblue')
plt.title('Distribution of Sentiment Scores')
plt.xlabel('Sentiment Score')
plt.ylabel('Frequency')
plt.show()

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

# Add a title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "British Airways Review Analysis"
subtitle.text = "Summary of Customer Feedback and Insights"

# Add a content slide for visualizations
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Analysis Overview"

# Add Word Cloud
slide.shapes.add_picture('word_cloud.png', Inches(0.5), Inches(1.5), width=Inches(4))

# Add Average Ratings Bar Chart
slide.shapes.add_picture('average_ratings.png', Inches(5), Inches(1.5), width=Inches(4))

# Add Sentiment Distribution Histogram
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Sentiment Analysis and Recommendations"

slide.shapes.add_picture('sentiment_distribution.png', Inches(0.5), Inches(1.5), width=Inches(4))

# Add Recommendation Rate Pie Chart
slide.shapes.add_picture('recommendation_rate.png', Inches(5), Inches(1.5), width=Inches(4))

# Add Conclusions Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusions and Recommendations"

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(3))
text_frame = text_box.text_frame
text_frame.text = (
    "1. **Strengths**: High ratings for Cabin Staff Service and Seat Comfort.\n"
    "2. **Areas for Improvement**: Wifi & Connectivity and Value for Money could be improved.\n"
    "3. **Recommendation Rate**: Approximately XX% of customers recommend British Airways.\n"
    "4. **Sentiment Analysis**: Average sentiment score indicates positive feedback with some neutral comments.\n"
    "5. **Word Cloud Insights**: Common keywords highlight customer focus on service, comfort, and overall value."
)

prs.save("British_Airways_Review_Analysis.pptx")
